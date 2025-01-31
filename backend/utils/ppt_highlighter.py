from pptx import Presentation
from pptx.dml.color import RGBColor
import logging
import json
from pathlib import Path

logger = logging.getLogger(__name__)


def get_highlight_color(color_name: str) -> tuple:
    """Get RGB color tuple for a given color name"""
    colors = {
        "green": (0, 204, 0),  # Green for main ideas
        "yellow": (255, 255, 0),  # Yellow for vocabulary
        "pink": (255, 179, 179),  # Pink for questions
        "blue": (128, 128, 255),  # Lighter blue for sub-ideas
    }
    return colors.get(color_name, (255, 255, 0))  # Default to yellow


def add_highlights(highlight_data: list, filename: str) -> str:
    """
    Add color-coded highlights to a PowerPoint file.

    Args:
        highlight_data (list): List of dictionaries containing content and highlight information
        filename (str): Path to the PPTX file

    Returns:
        str: Path to the highlighted PPTX file
    """
    logger.info("Adding highlights to %s", filename)

    try:
        # Ensure input file exists
        input_path = Path(filename)
        if not input_path.exists():
            logger.error("Input file not found: %s", filename)
            return None

        print("input_path :", input_path)

        # Open presentation
        prs = Presentation(str(input_path))

        # Create output directory if it doesn't exist
        output_dir = Path("files")
        output_dir.mkdir(exist_ok=True)

        # Process each highlight
        for highlight in highlight_data:
            content = highlight["content"]
            color = get_highlight_color(highlight["highlight_type"]["color"])

            # Search through all slides and shapes
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text_frame"):
                        for paragraph in shape.text_frame.paragraphs:
                            if content in paragraph.text:
                                # Split text and apply highlighting
                                for run in paragraph.runs:
                                    if content in run.text:
                                        run.font.fill.solid()
                                        run.font.fill.fore_color.rgb = RGBColor(*color)
                                        # Add comment as speaker notes if needed
                                        if not slide.has_notes_slide:
                                            slide.notes_slide
                                        slide.notes_slide.notes_text_frame.text += (
                                            f"\nHighlight: {highlight['explanation']}"
                                        )

        # Save highlighted presentation
        output_path = str(output_dir / f"highlighted_{input_path.name}")
        prs.save(output_path)
        logger.info("Saved highlighted file to %s", output_path)
        return output_path

    except Exception as e:
        logger.error("Error highlighting PPTX: %s", str(e))
        return None
